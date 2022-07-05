import collections 
import collections.abc
from pptx import Presentation
import os

#for images
from pptx.util import Inches

#for shapes

#from pptx.enum.shapes import MSO_SHAPE
#from pptx.enum.shapes import MSO_SHAPE_TYPE
#from pptx.enum.dml import MSO_SHAPE_COLOR

#For graphs - slide 5

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_TICK_MARK
from pptx.util import Pt


#slide 1

#creates ppt presentation to add slides to
prs = Presentation()

#refers to layouts of each slide
#0=title slide 1=title and content 3=section header
slide1_register = prs.slide_layouts[0]

#add initial slide to presentation
slide1 = prs.slides.add_slide(slide1_register)

#main top placeholder
title1 = slide1.shapes.title

#placeholder = item in layout 
subtitle1 = slide1.placeholders[1]

#insert text
title1.text = "Hello, World!"
subtitle1.text = "Myu was here!"

#drugi slajd - dodajemy bulletpoity

#creating bullet point slide
slide2_register = prs.slide_layouts[1]

#adding slide
slide2 = prs.slides.add_slide(slide2_register)

#edit bullet point slide

title2 = slide2.shapes.title
title2.text = 'Adding a Bullet Slide by Myu'



bullet_point_box = slide2.shapes

bullets_lvl_1 = bullet_point_box.placeholders[1]
bullets_lvl_1.text = "Pierwszy stopien"



bullets_lvl_2 = bullets_lvl_1.text_frame.add_paragraph()
bullets_lvl_2.text = 'Drugi stopien'
bullets_lvl_2.level = 1


bullets_lvl_2 = bullets_lvl_1.text_frame.add_paragraph()
bullets_lvl_2.text = 'Drugi stopien kolejna linia'
bullets_lvl_2.level = 1

bullets_lvl_3 = bullets_lvl_1.text_frame.add_paragraph()
bullets_lvl_3.text = "Trutututu"
bullets_lvl_3.level = 2

bullets_lvl_4 = bullets_lvl_1.text_frame.add_paragraph()
bullets_lvl_4.text = "Plum"
bullets_lvl_4.level = 3

#slide 3
#create bullet point and picture slide

slide3_register = prs.slide_layouts[5]
slide3 = prs.slides.add_slide(slide3_register)

title3 = slide3.shapes.title
title3.text = "Picture Time!"

#add image
img1 = 'zwyciestwo.jpeg'

from_left = Inches(0)
from_top = Inches(1.5)

add_picture = slide3.shapes.add_picture(img1,from_left,from_top)


# slide 4  - autoshapes

slide4_register = prs.slide_layouts[5]
slide4 = prs.slides.add_slide(slide4_register)

title4 = slide4.shapes.title
title4.text = "Shapework"

#create shapes - we need to import subsection from pptx module 

#Shape 1

left1 = top1 = width1 = height1 = Inches(1)
#add_shape1 = slide4.shapes.add_shape MSO_SHAPE.ROUNDED_RECTANGLE,left1,top1,width1,height1

# Shape 2 

left2 = Inches(6)
top2 = Inches(2)
width2 = height2 = Inches (2)
#arrow1 = slide4.shapes.add_shape MSO_SHAPE.DOWN_ARROW,left2,top2,width2,height2clear

#edit arrow - colour"

#fill_arrow1 = arrow1.fill
#fill_arrow1.solid()
#fill_arrow1.fore_color.theme_color = MSO_THEME_COLOR

#edit arrow - rotate
#arrow1.rotation = 90

# Slide 5 - Graph slide 
slide5_register = prs.slide_layouts[5]
slide5 = prs.slides.add_slide(slide5_register)

title5 = slide5.shapes.title
title5.text = "Let's make a Graph"

#Bulit graph
graph_info = CategoryChartData()
graph_info.categories = ["A", "B", "C"]
graph_info.add_series("Series 1", (15, 11, 18))

#Add graph to Siles with Positioning
left_graph = Inches(2)
top_graph = Inches(2)
width_graph = Inches(6)
height_graph = Inches(4)
graph1_frame = slide5.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left_graph, top_graph, width_graph, height_graph, graph_info)


graph1 = graph1_frame.chart

#edit graph
category_axis = graph1.category_axis
category_axis.has_major_gridlines = True
category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
category_axis.tick_labels.font.italic = True
category_axis.tick_labels.font.size = Pt(24)



#generowanie slajdow

prs.save('test2.pptx')
